from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Цвета бренда ────────────────────────────────────────────────────────
BG     = RGBColor(0x0d, 0x11, 0x17)
CARD   = RGBColor(0x16, 0x1b, 0x22)
BORDER = RGBColor(0x30, 0x36, 0x3d)
FG     = RGBColor(0xe6, 0xed, 0xf3)
MUTED  = RGBColor(0x8b, 0x94, 0x9e)
ACCENT = RGBColor(0x58, 0xa6, 0xff)
GREEN  = RGBColor(0x3f, 0xb9, 0x50)
WARN   = RGBColor(0xe3, 0xb3, 0x41)
DANGER = RGBColor(0xf8, 0x51, 0x49)
PURPLE = RGBColor(0xbc, 0x8c, 0xff)
WHITE  = RGBColor(0xff, 0xff, 0xff)

# ── Размер слайда 16:9 ──────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

blank = prs.slide_layouts[6]  # пустой макет

# ════════════════════════════════════════════════════════════════════════
# Вспомогательные функции
# ════════════════════════════════════════════════════════════════════════

def new_slide():
    sl = prs.slides.add_slide(blank)
    # Тёмный фон
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return sl


def add_rect(sl, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Emu
    shape = sl.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(sl, x, y, w, h, text, size=Pt(14), bold=False, color=FG,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = sl.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, size=Pt(13), bold=False, color=FG,
             align=PP_ALIGN.LEFT, space_before=Pt(0), italic=False):
    p   = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def tag_box(sl, x, y, text, color=ACCENT, bg_alpha=None):
    """Маленький тег-пилюля"""
    w = Inches(2.2)
    h = Inches(0.3)
    rect = add_rect(sl, x, y, w, h, fill_color=CARD, line_color=color, line_width=Pt(0.75))
    tf = rect.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.size  = Pt(8)
    run.font.bold  = True
    run.font.color.rgb = color
    return rect


def card(sl, x, y, w, h, title=None, lines=None, icon=None,
         fill=CARD, border=BORDER, title_color=FG, line_color=MUTED,
         title_size=Pt(13)):
    rect = add_rect(sl, x, y, w, h, fill_color=fill, line_color=border, line_width=Pt(0.75))
    pad  = Inches(0.18)
    cy   = y + pad
    if icon:
        add_textbox(sl, x+pad, cy, w-2*pad, Inches(0.35), icon, size=Pt(18))
        cy += Inches(0.38)
    if title:
        add_textbox(sl, x+pad, cy, w-2*pad, Inches(0.32), title,
                    size=title_size, bold=True, color=title_color)
        cy += Inches(0.35)
    if lines:
        for ln in lines:
            add_textbox(sl, x+pad, cy, w-2*pad, Inches(0.28), ln,
                        size=Pt(10.5), color=line_color, wrap=True)
            cy += Inches(0.27)
    return rect


def checklist(sl, x, y, w, items, size=Pt(11.5), spacing=Inches(0.32)):
    cy = y
    for item in items:
        add_textbox(sl, x, cy, Inches(0.28), spacing, "✓",
                    size=size, bold=True, color=GREEN)
        add_textbox(sl, x + Inches(0.3), cy, w - Inches(0.3), spacing + Inches(0.04),
                    item, size=size, color=FG, wrap=True)
        cy += spacing
    return cy


def stat_box(sl, x, y, w, h, num, lbl, num_color=ACCENT):
    rect = add_rect(sl, x, y, w, h, fill_color=CARD, line_color=BORDER, line_width=Pt(0.75))
    add_textbox(sl, x, y + Inches(0.15), w, Inches(0.6), num,
                size=Pt(32), bold=True, color=num_color, align=PP_ALIGN.CENTER)
    add_textbox(sl, x, y + Inches(0.75), w, Inches(0.35), lbl,
                size=Pt(9), color=MUTED, align=PP_ALIGN.CENTER)
    return rect


def section_slide(num_str, tag_text, tag_color, title, subtitle, emoji):
    sl = new_slide()
    # Лёгкое свечение по центру
    glow = add_rect(sl, Inches(3.5), Inches(1.5), Inches(6.33), Inches(4.5),
                    fill_color=RGBColor(0x0d, 0x15, 0x22), line_color=None)

    cx = SW / 2
    # Emoji
    add_textbox(sl, cx - Inches(1.5), Inches(1.2), Inches(3), Inches(1),
                emoji, size=Pt(54), align=PP_ALIGN.CENTER, color=FG)
    # Тег
    tag_w = Inches(2.5)
    tag_box(sl, cx - tag_w/2, Inches(2.35), tag_text, color=tag_color)
    # Заголовок
    add_textbox(sl, Inches(1), Inches(2.75), Inches(11.33), Inches(1.2),
                title, size=Pt(44), bold=True, color=FG, align=PP_ALIGN.CENTER)
    # Подзаголовок
    add_textbox(sl, Inches(1.5), Inches(3.95), Inches(10.33), Inches(0.7),
                subtitle, size=Pt(17), color=MUTED, align=PP_ALIGN.CENTER)
    # Номер слайда (маленький)
    add_textbox(sl, Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.3),
                num_str, size=Pt(9), color=BORDER)
    return sl


def slide_num(sl, n):
    add_textbox(sl, Inches(12.8), Inches(7.1), Inches(0.5), Inches(0.3),
                str(n), size=Pt(9), color=BORDER)


# ════════════════════════════════════════════════════════════════════════
# Слайд 1 — Титул
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()

# Лого WeldTeam MES
txb = sl.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.33), Inches(1.5))
tf  = txb.text_frame
p   = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r1 = p.add_run(); r1.text = "Weld";          r1.font.size = Pt(56); r1.font.bold = True;  r1.font.color.rgb = FG
r2 = p.add_run(); r2.text = "Team";          r2.font.size = Pt(56); r2.font.bold = True;  r2.font.color.rgb = ACCENT
r3 = p.add_run(); r3.text = "  MES";         r3.font.size = Pt(44); r3.font.bold = False; r3.font.color.rgb = MUTED

add_textbox(sl, Inches(1), Inches(3.15), Inches(11.33), Inches(0.55),
            "Система управления сварочным производством",
            size=Pt(19), color=MUTED, align=PP_ALIGN.CENTER)

# Зелёная строка статуса
status_rect = add_rect(sl, Inches(3.8), Inches(3.85), Inches(5.73), Inches(0.45),
                       fill_color=RGBColor(0x0d, 0x1e, 0x10),
                       line_color=RGBColor(0x3f, 0xb9, 0x50), line_width=Pt(0.75))
add_textbox(sl, Inches(3.8), Inches(3.87), Inches(5.73), Inches(0.4),
            "● Система работает  ·  Данные внесены  ·  Фундамент заложен",
            size=Pt(11), color=GREEN, align=PP_ALIGN.CENTER)

# Три тега
tags = [("Для отдела ИТО", GREEN), ("Для отдела ОТК", ACCENT), ("Цех BS · GWM", PURPLE)]
tx = Inches(3.0)
for t_text, t_color in tags:
    tw = Inches(2.3)
    tr = add_rect(sl, tx, Inches(4.55), tw, Inches(0.38),
                  fill_color=CARD, line_color=t_color, line_width=Pt(0.75))
    add_textbox(sl, tx, Inches(4.57), tw, Inches(0.35),
                t_text.upper(), size=Pt(8.5), bold=True, color=t_color, align=PP_ALIGN.CENTER)
    tx += Inches(2.5)

slide_num(sl, "1 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 2 — Фундамент (статистика)
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(0.5), Inches(0.4), "Фундамент", color=GREEN)
add_textbox(sl, Inches(0.5), Inches(0.8), Inches(12), Inches(0.65),
            "Что уже реализовано и работает",
            size=Pt(30), bold=True, color=FG)

stats_data = [
    ("569",    "Сварочных\nпистолетов", ACCENT),
    ("16 680", "Сварочных\nточек",      PURPLE),
    ("6",      "Моделей\nавто",         WARN),
    ("2 606",  "Записей\nплана ТО",     GREEN),
    ("3",      "Произв.\nлинии",        DANGER),
]
sb_w = Inches(2.2)
sb_h = Inches(1.4)
sx   = Inches(0.5)
for num, lbl, col in stats_data:
    stat_box(sl, sx, Inches(1.65), sb_w, sb_h, num, lbl, col)
    sx += Inches(2.5)

# Хайлайт
hl = add_rect(sl, Inches(0.5), Inches(3.25), Inches(12.33), Inches(1.0),
              fill_color=RGBColor(0x0e, 0x18, 0x28),
              line_color=RGBColor(0x1d, 0x40, 0x6b), line_width=Pt(0.75))
add_rect(sl, Inches(0.5), Inches(3.25), Inches(0.07), Inches(1.0),
         fill_color=ACCENT, line_color=None)
add_textbox(sl, Inches(0.75), Inches(3.35), Inches(11.9), Inches(0.85),
            "Вся структура цеха занесена в базу данных: бренды → станции → пистолеты → точки сварки → уставки "
            "техпроцесса. Это фундаментальная работа, которая позволяет строить любой аналитический и операционный функционал.",
            size=Pt(12.5), color=FG, wrap=True)

# Дополнительный блок — архитектура данных
add_textbox(sl, Inches(0.5), Inches(4.45), Inches(12), Inches(0.35),
            "Иерархия данных:", size=Pt(11), bold=True, color=MUTED)
flow_items = ["Бренд (линия)", "Станция", "Пистолет", "Точки сварки", "Уставки техпроцесса", "История изменений"]
fx = Inches(0.5)
for i, fi in enumerate(flow_items):
    fw = Inches(1.9)
    fc = ACCENT if i % 2 == 0 else MUTED
    add_rect(sl, fx, Inches(4.85), fw, Inches(0.38),
             fill_color=CARD, line_color=fc, line_width=Pt(0.5))
    add_textbox(sl, fx, Inches(4.87), fw, Inches(0.35), fi,
                size=Pt(9.5), color=fc, align=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_textbox(sl, fx + fw, Inches(4.9), Inches(0.25), Inches(0.3),
                    "→", size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)
    fx += Inches(2.18)

slide_num(sl, "2 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 3 — Технология
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(0.5), Inches(0.4), "Технология", color=ACCENT)
add_textbox(sl, Inches(0.5), Inches(0.8), Inches(12), Inches(0.65),
            "Как работает система", size=Pt(30), bold=True, color=FG)

tech_cards = [
    ("🌐", "Веб-приложение",  ["Открывается в любом браузере", "Не требует установки", "Работает на ПК, телефоне", "и планшете в сети цеха"]),
    ("🗄️", "База данных",     ["Централизованное хранение", "Актуальные данные в", "реальном времени для всех", "сотрудников"]),
    ("📱", "Мобильная версия",["Страницы адаптированы", "под телефон", "Регистрация ТО и дефектов", "прямо у станка"]),
    ("🔒", "Аудит изменений", ["Каждое изменение уставок", "записывается с датой", "и комментарием", "История не удаляется"]),
    ("⚡", "Валидация данных", ["Проверка замеров ТО:", "допуск ±100 А по току,", "±50 daN по давлению", "Выход — предупреждение"]),
    ("🌙", "Тёмная тема",     ["Комфортная работа", "в любых условиях", "освещения цеха", "Сохраняется между сессиями"]),
]

cw = Inches(3.9)
ch = Inches(2.3)
for i, (ico, ttl, lns) in enumerate(tech_cards):
    cx = Inches(0.5) + (i % 3) * Inches(4.27)
    cy = Inches(1.6) + (i // 3) * Inches(2.5)
    card(sl, cx, cy, cw, ch, title=ttl, icon=ico, lines=lns)

slide_num(sl, "3 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 4 — Секция ИТО
# ════════════════════════════════════════════════════════════════════════
sl = section_slide("4 / 15", "Отдел ИТО", GREEN,
                   "Техническое обслуживание",
                   "Планирование, выполнение и контроль ТО сварочных пистолетов",
                   "🔧")


# ════════════════════════════════════════════════════════════════════════
# Слайд 5 — ТО Plan
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(0.5), Inches(0.4), "ИТО · Планирование", color=GREEN)
add_textbox(sl, Inches(0.5), Inches(0.8), Inches(12), Inches(0.65),
            "Ежедневный план ТО", size=Pt(30), bold=True, color=FG)

to_items = [
    "Годовой план ТО загружен — 2 606 записей по всем брендам",
    "Назначить пистолеты на ТО одним кликом — с фильтром по бренду и дате",
    "В задаче видно: номер и тип пистолета, бренд, станция",
    "Взять задачу в работу прямо с экрана телефона у станка",
    "Записать ТО из задачи — задача автоматически закрывается",
    "Прогресс выполнения плана по брендам с % готовности",
]
checklist(sl, Inches(0.5), Inches(1.65), Inches(7.2), to_items, size=Pt(12))

# Карточка потока работ
flow_card = add_rect(sl, Inches(8.3), Inches(1.55), Inches(4.5), Inches(5.0),
                     fill_color=CARD, line_color=BORDER, line_width=Pt(0.75))
add_textbox(sl, Inches(8.5), Inches(1.7), Inches(4.1), Inches(0.3),
            "ПОТОК РАБОТЫ", size=Pt(8), bold=True, color=MUTED)

flow_steps = [
    ("📋 Назначить пистолеты", MUTED),
    ("🔧 Взять в работу",      ACCENT),
    ("📏 Снять замеры (3×ток, 3×давл.)", WARN),
    ("✅ ТО зафиксировано",    GREEN),
]
fy = Inches(2.1)
for i, (step, col) in enumerate(flow_steps):
    step_rect = add_rect(sl, Inches(8.5), fy, Inches(4.1), Inches(0.48),
                         fill_color=BG, line_color=col, line_width=Pt(0.75))
    add_textbox(sl, Inches(8.5), fy + Inches(0.02), Inches(4.1), Inches(0.44),
                step, size=Pt(11), bold=True, color=col, align=PP_ALIGN.CENTER)
    fy += Inches(0.52)
    if i < len(flow_steps) - 1:
        add_textbox(sl, Inches(10.3), fy - Inches(0.1), Inches(0.5), Inches(0.3),
                    "↓", size=Pt(14), color=MUTED, align=PP_ALIGN.CENTER)
        fy += Inches(0.2)

slide_num(sl, "5 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 6 — ТО Уставки
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(0.5), Inches(0.4), "ИТО · Уставки", color=WARN)
add_textbox(sl, Inches(0.5), Inches(0.8), Inches(12), Inches(0.65),
            "Управление параметрами сварки", size=Pt(30), bold=True, color=FG)

param_items = [
    "Уставки хранятся в БД для каждого пистолета и каждой точки",
    "Все параметры: давление (daN), к.трансформации, ток 1/2, времена циклограммы",
    "Изменение уставок — с обязательным комментарием «зачем меняем»",
    "Полный архив: когда, какие параметры, какой комментарий",
    "Старые уставки не удаляются — сохраняется полная история",
    "Привязка замеров ТО к конкретной уставке (факт/план)",
]
checklist(sl, Inches(0.5), Inches(1.65), Inches(7.5), param_items, size=Pt(12))

# Карточка — что в архиве
arc = add_rect(sl, Inches(8.3), Inches(1.55), Inches(4.5), Inches(4.0),
               fill_color=CARD, line_color=BORDER, line_width=Pt(0.75))
add_textbox(sl, Inches(8.5), Inches(1.7), Inches(4.1), Inches(0.35),
            "📊  Что видно в архиве уставок", size=Pt(12), bold=True, color=FG)
archive_lines = [
    ("Дата введения уставки",       MUTED, False),
    ("Режим сварки",                MUTED, False),
    ("Давление, К.Трансформации",   MUTED, False),
    ("Ток 1 и Ток 2",              MUTED, False),
    ("Комментарий к изменению ← новое", GREEN, True),
    ("Активна / в архиве",         MUTED, False),
]
ay = Inches(2.15)
for aline, acol, abold in archive_lines:
    add_textbox(sl, Inches(8.65), ay, Inches(0.18), Inches(0.28), "·",
                size=Pt(12), color=acol)
    add_textbox(sl, Inches(8.88), ay, Inches(3.8), Inches(0.3), aline,
                size=Pt(11), color=acol, bold=abold)
    ay += Inches(0.32)

slide_num(sl, "6 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 7 — Секция ОТК
# ════════════════════════════════════════════════════════════════════════
sl = section_slide("7 / 15", "Отдел ОТК", ACCENT,
                   "Управление дефектами",
                   "Регистрация, расследование и закрытие дефектов сварки",
                   "🔍")


# ════════════════════════════════════════════════════════════════════════
# Слайд 8 — Дефекты
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(0.5), Inches(0.4), "ОТК · Дефекты", color=ACCENT)
add_textbox(sl, Inches(0.5), Inches(0.8), Inches(12), Inches(0.65),
            "Регистрация и закрытие дефектов", size=Pt(30), bold=True, color=FG)

def_items = [
    "16 кодов дефектов из техпроцесса (CR, SN, LP, BN, SW и др.)",
    "Регистрация: модель авто → номер точки → код дефекта",
    "В карточке видно: тип пистолета, бренд, станция — сразу найти",
    "Workflow: Зарегистрирован → Закрыть (без лишних шагов)",
    "При закрытии: первопричина + контрмера + ответственный",
    "История закрытых дефектов с поиском по периоду",
]
checklist(sl, Inches(0.5), Inches(1.65), Inches(7.3), def_items, size=Pt(12))

# Карточка кодов
codes_card = add_rect(sl, Inches(8.3), Inches(1.55), Inches(4.5), Inches(4.3),
                      fill_color=CARD, line_color=BORDER, line_width=Pt(0.75))
add_textbox(sl, Inches(8.5), Inches(1.7), Inches(4.1), Inches(0.3),
            "16 КОДОВ ДЕФЕКТОВ", size=Pt(8), bold=True, color=MUTED)

codes = ["CR", "SN", "LP", "BN", "SW", "BT", "MS", "ME", "IE", "MN", "NA", "EO", "BE", "P", "MI", "EMU"]
cx2 = Inches(8.5)
cy2 = Inches(2.1)
cw2 = Inches(0.72)
ch2 = Inches(0.33)
for i, code in enumerate(codes):
    if i > 0 and i % 5 == 0:
        cy2 += Inches(0.42)
        cx2 = Inches(8.5)
    c_rect = add_rect(sl, cx2, cy2, cw2, ch2,
                      fill_color=RGBColor(0x25, 0x0d, 0x0c),
                      line_color=DANGER, line_width=Pt(0.5))
    add_textbox(sl, cx2, cy2, cw2, ch2, code,
                size=Pt(9.5), bold=True, color=DANGER, align=PP_ALIGN.CENTER)
    cx2 += Inches(0.82)

add_textbox(sl, Inches(8.5), Inches(4.95), Inches(4.1), Inches(0.35),
            "Словарь кодов хранится в БД и может пополняться",
            size=Pt(9.5), color=MUTED, italic=True, wrap=True)

slide_num(sl, "8 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 9 — Аналитика
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(0.5), Inches(0.4), "ОТК · Аналитика", color=ACCENT)
add_textbox(sl, Inches(0.5), Inches(0.8), Inches(12), Inches(0.65),
            "Аналитика дефектов", size=Pt(30), bold=True, color=FG)

an_items = [
    "Фильтры: этот месяц / 3 мес / 6 мес / год / свой период",
    "Фильтр по модели автомобиля и конкретной станции",
    "Диаграмма: наглядно видно какие коды доминируют",
    "Таблица разбивки по станциям — где больше проблем",
    "Самый частый код дефекта за период — на главном экране",
]
checklist(sl, Inches(0.5), Inches(1.65), Inches(7.3), an_items, size=Pt(12))

# Хайлайт
hl2 = add_rect(sl, Inches(0.5), Inches(4.15), Inches(7.3), Inches(0.9),
               fill_color=RGBColor(0x0e, 0x18, 0x28),
               line_color=RGBColor(0x1d, 0x40, 0x6b), line_width=Pt(0.75))
add_rect(sl, Inches(0.5), Inches(4.15), Inches(0.07), Inches(0.9),
         fill_color=ACCENT, line_color=None)
add_textbox(sl, Inches(0.75), Inches(4.25), Inches(6.9), Inches(0.7),
            "Аналитика отвечает на вопрос: где и что чаще всего ломается?\n"
            "Позволяет принимать системные решения, а не тушить пожары",
            size=Pt(12), color=FG, wrap=True)

# Легенда диаграммы вместо реального пончика
donut_card = add_rect(sl, Inches(8.3), Inches(1.55), Inches(4.5), Inches(4.5),
                      fill_color=CARD, line_color=BORDER, line_width=Pt(0.75))
add_textbox(sl, Inches(8.5), Inches(1.7), Inches(4.1), Inches(0.3),
            "РАСПРЕДЕЛЕНИЕ ДЕФЕКТОВ", size=Pt(8), bold=True, color=MUTED)

legend = [
    (DANGER, "CR — Трещина",         "35%"),
    (WARN,   "SN — Малое ядро",      "20%"),
    (ACCENT, "LP — Непровар",        "17%"),
    (GREEN,  "BN — Прожог",          "13%"),
    (PURPLE, "прочие коды...",        "15%"),
]
ly = Inches(2.15)
bar_x = Inches(8.5)
bar_total = Inches(4.1)
for lc, lt, lpct in legend:
    pct_val = int(lpct.replace('%', ''))
    bar_w   = bar_total * pct_val / 100
    add_rect(sl, bar_x, ly, bar_total, Inches(0.32),
             fill_color=RGBColor(0x20, 0x28, 0x35), line_color=None)
    add_rect(sl, bar_x, ly, bar_w, Inches(0.32),
             fill_color=lc, line_color=None)
    add_textbox(sl, bar_x + Inches(0.1), ly + Inches(0.02), bar_total - Inches(0.15), Inches(0.28),
                f"{lt}  {lpct}", size=Pt(9.5), bold=False, color=FG)
    ly += Inches(0.55)

slide_num(sl, "9 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 10 — Обзор данных (Explorer)
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(0.5), Inches(0.4), "ИТО + ОТК · Обзор данных", color=PURPLE)
add_textbox(sl, Inches(0.5), Inches(0.8), Inches(12), Inches(0.65),
            "Полный паспорт любого объекта цеха", size=Pt(30), bold=True, color=FG)

exp_cards = [
    ("🔫", "По пистолету", ACCENT, [
        "Паспорт: номер, тип, бренд, станция",
        "Активные уставки по каждой точке",
        "Архив уставок с комментариями",
        "История всех ТО",
        "Список точек которые он варит",
        "Все дефекты пистолета",
    ]),
    ("🏭", "По станции", WARN, [
        "Все пистолеты с нумерацией",
        "Производимые модели и кол-во точек",
        "Сводка дефектов по кодам",
        "Последние ТО на станции",
    ]),
    ("🎯", "По точке сварки", GREEN, [
        "Активный пистолет и уставка",
        "История: когда какой пистолет варил",
        "Все дефекты на этой точке",
        "Последние ТО пистолета",
    ]),
]
ecw = Inches(4.0)
ech = Inches(4.2)
for i, (ico, ttl, tcol, lns) in enumerate(exp_cards):
    ex = Inches(0.5) + i * Inches(4.4)
    card(sl, ex, Inches(1.6), ecw, ech, title=ttl, icon=ico, lines=lns,
         title_color=tcol, border=tcol)

# Хайлайт внизу
hl3 = add_rect(sl, Inches(0.5), Inches(5.95), Inches(12.33), Inches(0.85),
               fill_color=RGBColor(0x0e, 0x18, 0x28),
               line_color=RGBColor(0x1d, 0x40, 0x6b), line_width=Pt(0.75))
add_rect(sl, Inches(0.5), Inches(5.95), Inches(0.07), Inches(0.85),
         fill_color=ACCENT, line_color=None)
add_textbox(sl, Inches(0.75), Inches(6.05), Inches(11.9), Inches(0.7),
            "Ключевое преимущество: любой инженер за 10 секунд получает полную картину по любому элементу цеха — "
            "не нужно искать в Excel, звонить коллегам или поднимать бумажные журналы",
            size=Pt(11.5), color=FG, wrap=True)

slide_num(sl, "10 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 11 — Секция Будущее
# ════════════════════════════════════════════════════════════════════════
sl = section_slide("11 / 15", "Дальнейшее развитие", PURPLE,
                   "Что планируется дальше",
                   "Фундамент заложен — теперь надстраиваем",
                   "🚀")


# ════════════════════════════════════════════════════════════════════════
# Слайд 12 — Weld Balance → А-лист (проблема и решение)
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(0.5), Inches(0.35), "ИТО · Следующий большой шаг", color=WARN)
add_textbox(sl, Inches(0.5), Inches(0.75), Inches(12.33), Inches(0.72),
            "Weld Balance → А-лист: один источник данных",
            size=Pt(28), bold=True, color=FG)

# Подзаголовки колонок
add_textbox(sl, Inches(0.5), Inches(1.55), Inches(5.9), Inches(0.35),
            "КАК СЕЙЧАС", size=Pt(9), bold=True, color=DANGER)
add_rect(sl, Inches(0.5), Inches(1.88), Inches(5.9), Inches(0.03),
         fill_color=DANGER, line_color=None)

add_textbox(sl, Inches(7.0), Inches(1.55), Inches(5.9), Inches(0.35),
            "КАК БУДЕТ", size=Pt(9), bold=True, color=GREEN)
add_rect(sl, Inches(7.0), Inches(1.88), Inches(5.9), Inches(0.03),
         fill_color=GREEN, line_color=None)

# Разделитель
add_rect(sl, Inches(6.55), Inches(1.5), Inches(0.04), Inches(5.3),
         fill_color=BORDER, line_color=None)

# ЛЕВАЯ колонка — боль
pain_items = [
    ("📄", DANGER, "Weld Balance заполняется в Excel"),
    ("📋", DANGER, "А-лист заполняется ОТДЕЛЬНО — та же работа дважды"),
    ("✏️",  DANGER, "Изменение WB → нужно вручную обновить А-лист"),
    ("⏰", DANGER, "На станциях могут лежать устаревшие А-листы"),
    ("📞", DANGER, "ОТК и производство не знают об изменениях"),
]
py = Inches(2.05)
for ico, col, txt in pain_items:
    add_textbox(sl, Inches(0.5), py, Inches(0.35), Inches(0.38), ico, size=Pt(14))
    add_rect(sl, Inches(0.5), py, Inches(5.9), Inches(0.38),
             fill_color=RGBColor(0x1e, 0x0d, 0x0d), line_color=None)
    add_textbox(sl, Inches(0.9), py + Inches(0.05), Inches(5.4), Inches(0.3),
                txt, size=Pt(11), color=col, wrap=True)
    py += Inches(0.5)

# ПРАВАЯ колонка — решение
sol_items = [
    ("🗄️",  GREEN,  "Weld Balance ведётся в WeldTeam — один раз"),
    ("⚡",  GREEN,  "А-лист генерируется АВТОМАТИЧЕСКИ из БД"),
    ("🔄",  GREEN,  "Изменение WB → А-лист обновляется автоматически"),
    ("📧",  GREEN,  "Авторассылка уведомлений всем отделам"),
    ("🖨️",  GREEN,  "Печать нового А-листа — одна кнопка, работник относит на станцию"),
]
sy = Inches(2.05)
for ico, col, txt in sol_items:
    add_rect(sl, Inches(7.0), sy, Inches(5.9), Inches(0.38),
             fill_color=RGBColor(0x0d, 0x1e, 0x10), line_color=None)
    add_textbox(sl, Inches(7.0), sy, Inches(0.35), Inches(0.38), ico, size=Pt(14))
    add_textbox(sl, Inches(7.4), sy + Inches(0.05), Inches(5.4), Inches(0.3),
                txt, size=Pt(11), color=col, wrap=True)
    sy += Inches(0.5)

# Хайлайт внизу
hl_wb = add_rect(sl, Inches(0.5), Inches(4.75), Inches(12.33), Inches(0.75),
                 fill_color=RGBColor(0x1a, 0x16, 0x05),
                 line_color=WARN, line_width=Pt(0.75))
add_textbox(sl, Inches(0.75), Inches(4.85), Inches(11.9), Inches(0.6),
            "А-лист — опорный документ всего техпроцесса. Его используют ИТО, ОТК и производство. "
            "Автогенерация из WeldTeam гарантирует, что на станциях ВСЕГДА лежит актуальная документация.",
            size=Pt(11.5), color=FG, bold=False, wrap=True)

slide_num(sl, "12 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 13 — Roadmap (обновлённый)
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(0.5), Inches(0.35), "Roadmap", color=PURPLE)
add_textbox(sl, Inches(0.5), Inches(0.75), Inches(12), Inches(0.6),
            "Три фазы развития системы", size=Pt(28), bold=True, color=FG)

def roadmap_phase(sl, x, y, w, phase_lbl, phase_col, phase_bg, title, items):
    # Заголовок фазы
    add_rect(sl, x, y, w, Inches(0.48),
             fill_color=phase_bg, line_color=phase_col, line_width=Pt(0.75))
    txb = sl.shapes.add_textbox(x + Inches(0.15), y + Inches(0.04), w - Inches(0.3), Inches(0.4))
    tf = txb.text_frame
    p  = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = phase_lbl + "  "; r1.font.size = Pt(10); r1.font.bold = True; r1.font.color.rgb = phase_col
    r2 = p.add_run(); r2.text = title;             r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = FG
    # Строки
    iy = y + Inches(0.56)
    for item_ico, item_txt in items:
        add_rect(sl, x, iy, w, Inches(0.44),
                 fill_color=CARD, line_color=BORDER, line_width=Pt(0.4))
        add_textbox(sl, x + Inches(0.12), iy + Inches(0.02), Inches(0.3), Inches(0.38),
                    item_ico, size=Pt(13))
        add_textbox(sl, x + Inches(0.48), iy + Inches(0.07), w - Inches(0.6), Inches(0.3),
                    item_txt, size=Pt(10.5), color=FG, wrap=True)
        iy += Inches(0.5)

pw = Inches(4.0)
roadmap_phase(sl, Inches(0.5), Inches(1.5), pw,
              "ФАЗА 1", ACCENT, RGBColor(0x0e, 0x18, 0x28),
              "ИТО — Weld Balance → А-лист",
              [
                  ("📝", "Цифровое ведение Weld Balance в системе"),
                  ("⚡", "Автогенерация А-листа из данных БД"),
                  ("🔄", "Управление изменениями: версионирование"),
                  ("📧", "Авторассылка уведомлений при изменениях"),
                  ("🖨️", "Печать А-листа одной кнопкой"),
              ])

roadmap_phase(sl, Inches(4.67), Inches(1.5), pw,
              "ФАЗА 2", WARN, RGBColor(0x28, 0x1e, 0x0a),
              "ОТК — Ежедневные проверки",
              [
                  ("📋", "Технические карты ежедневных проверок"),
                  ("👷", "Ежедневный план для каждого работника ОТК"),
                  ("✅", "Отметка выполненных точек прямо у станка"),
                  ("📊", "Аналитика выполнения плана проверок"),
                  ("🔗", "Привязка нарушений к дефектам в системе"),
              ])

roadmap_phase(sl, Inches(8.84), Inches(1.5), pw,
              "ФАЗА 3", PURPLE, RGBColor(0x1a, 0x10, 0x28),
              "Полная автоматизация",
              [
                  ("🔒", "Авторизация и роли: ИТО, ОТК, Производство"),
                  ("🤖", "Автосчитывание данных с контроллеров"),
                  ("📈", "Тренды деградации и прогноз ТО"),
                  ("⚠️",  "Уведомления о просроченных ТО"),
                  ("📤", "Экспорт отчётов в Excel / PDF"),
              ])

slide_num(sl, "13 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 14 — ОТК: будущий workflow ежедневных проверок
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(0.5), Inches(0.35), "ОТК · Фаза 2 — план проверок", color=WARN)
add_textbox(sl, Inches(0.5), Inches(0.75), Inches(12.33), Inches(0.65),
            "Ежедневные операционные проверки ОТК",
            size=Pt(28), bold=True, color=FG)

# Левая — суть
add_textbox(sl, Inches(0.5), Inches(1.55), Inches(6.3), Inches(0.3),
            "Что это и зачем", size=Pt(11), bold=True, color=MUTED)
otk_items = [
    "Каждый работник ОТК получает ежедневный план: какие точки проверить",
    "Технические карты содержат: точка, модель, параметры допуска, метод проверки",
    "Работник открывает план на телефоне → идёт к станции → отмечает результат",
    "Система знает: выполнено ли 100% плана за смену или остались пропуски",
    "При выявлении нарушения — регистрируется дефект прямо из проверки",
    "Аналитика: какие точки чаще всего нарушают параметры",
]
checklist(sl, Inches(0.5), Inches(1.95), Inches(6.3), otk_items, size=Pt(11.5))

# Правая — поток работы
flow_card = add_rect(sl, Inches(7.1), Inches(1.45), Inches(5.7), Inches(5.3),
                     fill_color=CARD, line_color=BORDER, line_width=Pt(0.75))
add_textbox(sl, Inches(7.3), Inches(1.6), Inches(5.3), Inches(0.3),
            "ПОТОК РАБОТЫ — СМЕНА ОТК", size=Pt(8), bold=True, color=MUTED)

otk_flow = [
    ("📋", ACCENT,  "А-лист формирует план\nпроверок на смену"),
    ("📱", WARN,    "Работник открывает план\nна телефоне / планшете"),
    ("🚶", FG,      "Идёт к станции, проверяет\nуказанные точки"),
    ("✅", GREEN,   "Отмечает результат:\nОК / нарушение"),
    ("🚨", DANGER,  "При нарушении — авто-\nрегистрация дефекта в системе"),
    ("📊", PURPLE,  "Аналитика за смену /\nмесяц в реальном времени"),
]
ffy = Inches(2.0)
for ffico, ffcol, fftxt in otk_flow:
    add_rect(sl, Inches(7.3), ffy, Inches(5.3), Inches(0.54),
             fill_color=BG, line_color=ffcol, line_width=Pt(0.6))
    add_textbox(sl, Inches(7.3), ffy, Inches(0.5), Inches(0.54), ffico, size=Pt(16), align=PP_ALIGN.CENTER)
    add_textbox(sl, Inches(7.85), ffy + Inches(0.05), Inches(4.6), Inches(0.44),
                fftxt, size=Pt(10.5), color=ffcol, wrap=True)
    ffy += Inches(0.62)

# Зависимость от А-листа
dep = add_rect(sl, Inches(0.5), Inches(5.65), Inches(12.33), Inches(0.65),
               fill_color=RGBColor(0x28, 0x1e, 0x0a),
               line_color=WARN, line_width=Pt(0.75))
add_textbox(sl, Inches(0.75), Inches(5.72), Inches(11.9), Inches(0.52),
            "⚡  Почему Фаза 2 идёт после Фазы 1: ОТК строит свою работу на А-листе. "
            "Сначала А-лист должен генерироваться автоматически — только тогда план проверок будет всегда актуальным.",
            size=Pt(11), color=WARN, wrap=True)

slide_num(sl, "14 / 15")


# ════════════════════════════════════════════════════════════════════════
# Слайд 15 — Итог
# ════════════════════════════════════════════════════════════════════════
sl = new_slide()
tag_box(sl, Inches(5.4), Inches(0.3), "Итог", color=GREEN)

txb = sl.shapes.add_textbox(Inches(1), Inches(0.7), Inches(11.33), Inches(1.1))
tf  = txb.text_frame
p   = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r1 = p.add_run(); r1.text = "Самое главное — ";  r1.font.size = Pt(40); r1.font.bold = True; r1.font.color.rgb = FG
r2 = p.add_run(); r2.text = "сделано";           r2.font.size = Pt(40); r2.font.bold = True; r2.font.color.rgb = ACCENT

add_textbox(sl, Inches(1.5), Inches(1.65), Inches(10.33), Inches(0.55),
            "База данных заполнена, все модули ИТО работают. "
            "Следующий шаг — Weld Balance → А-лист, затем ОТК.",
            size=Pt(13.5), color=MUTED, align=PP_ALIGN.CENTER, wrap=True)

sum_cards = [
    ("✅", "ИТО — сейчас",      GREEN,
     "ТО планируется и фиксируется. Уставки под контролем. Данные по каждому пистолету."),
    ("🔜", "ИТО — следующий шаг", WARN,
     "Weld Balance в системе → А-лист автоматически → рассылка и печать при изменениях."),
    ("🔜", "ОТК — после ИТО",   ACCENT,
     "Ежедневные планы проверок на базе А-листа. Фиксация нарушений. Аналитика смен."),
]
scw = Inches(3.9)
sch = Inches(2.1)
sx  = Inches(0.6)
for ico, ttl, tcol, descr in sum_cards:
    sr = add_rect(sl, sx, Inches(2.3), scw, sch,
                  fill_color=CARD, line_color=tcol, line_width=Pt(0.75))
    add_textbox(sl, sx + Inches(0.18), Inches(2.44), Inches(0.45), Inches(0.45), ico, size=Pt(20))
    add_textbox(sl, sx + Inches(0.18), Inches(2.88), scw - Inches(0.3), Inches(0.32),
                ttl, size=Pt(12.5), bold=True, color=tcol)
    add_textbox(sl, sx + Inches(0.18), Inches(3.24), scw - Inches(0.3), Inches(0.9),
                descr, size=Pt(10.5), color=MUTED, wrap=True)
    sx += Inches(4.27)

add_textbox(sl, Inches(1), Inches(4.6), Inches(11.33), Inches(0.55),
            "WeldTeam MES", size=Pt(24), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(sl, Inches(1), Inches(5.1), Inches(11.33), Inches(0.38),
            "Цех BS · GWM · 2026", size=Pt(13), color=MUTED, align=PP_ALIGN.CENTER)

slide_num(sl, "15 / 15")


# ════════════════════════════════════════════════════════════════════════
# Сохранение
# ════════════════════════════════════════════════════════════════════════
out = r"C:\Users\al.galimov\WeldTeam\BD\extracted_28_05\BD\web site\temp_server\WeldTeam_Presentation_v2.pptx"
prs.save(out)
print(f"Saved: {out}")
